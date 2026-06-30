"""Extracción de texto para la ingesta documental.

Soporta texto/CSV (decodificación directa), PDF (pypdf) y Word .docx (python-docx).
Para PDFs escaneados (sin capa de texto) intenta OCR si hay binario disponible
(pytesseract + pdf2image); si no, degrada de forma segura sin romper la carga.
"""
from __future__ import annotations

import io


def _pdf_text(raw: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception:  # pragma: no cover - dependencia ausente
        return ""
    try:
        reader = PdfReader(io.BytesIO(raw))
        parts = [(page.extract_text() or "") for page in reader.pages]
        text = "\n".join(parts).strip()
    except Exception:
        text = ""
    # PDF escaneado (poca/ninguna capa de texto): intenta OCR opcional.
    if len(text) < 20:
        ocr = _ocr_pdf(raw)
        if ocr:
            return ocr
    return text


def _ocr_pdf(raw: bytes) -> str:
    """OCR opcional: solo si pytesseract + pdf2image + tesseract están instalados."""
    try:  # pragma: no cover - depende de binario del sistema
        import pytesseract
        from pdf2image import convert_from_bytes

        images = convert_from_bytes(raw, dpi=200)
        return "\n".join(pytesseract.image_to_string(img, lang="spa+eng") for img in images).strip()
    except Exception:
        return ""


def _docx_text(raw: bytes) -> str:
    try:
        from docx import Document as Docx
    except Exception:  # pragma: no cover
        return ""
    try:
        d = Docx(io.BytesIO(raw))
        paras = [p.text for p in d.paragraphs if p.text]
        for table in d.tables:
            for row in table.rows:
                paras.append(" | ".join(c.text for c in row.cells))
        return "\n".join(paras).strip()
    except Exception:
        return ""


def _xlsx_text(raw: bytes) -> str:
    """Texto de un Excel .xlsx (todas las hojas, fila por fila)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    except Exception:  # pragma: no cover - dependencia/archivo inválido
        return ""
    out: list[str] = []
    try:
        for ws in wb.worksheets:
            out.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [("" if c is None else str(c)) for c in row]
                if any(cells):
                    out.append(" | ".join(cells))
    except Exception:
        return "\n".join(out).strip()
    return "\n".join(out).strip()


def _pptx_text(raw: bytes) -> str:
    """Texto de una presentación .pptx (títulos + viñetas de cada diapositiva)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
    except Exception:  # pragma: no cover
        return ""
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# Diapositiva {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
    return "\n".join(out).strip()


def _image_ocr(raw: bytes) -> str:
    """OCR opcional de una imagen (solo si pytesseract + tesseract están)."""
    try:  # pragma: no cover - depende de binario del sistema
        import pytesseract
        from PIL import Image
        return pytesseract.image_to_string(Image.open(io.BytesIO(raw)), lang="spa+eng").strip()
    except Exception:
        return ""


def _zip_text(raw: bytes, depth: int = 0) -> str:
    """Descomprime un .zip y extrae el texto de cada archivo soportado adentro."""
    if depth > 1:  # evita zips anidados profundos
        return ""
    import zipfile
    out: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            for info in z.infolist():
                if info.is_dir() or info.file_size > 50 * 1024 * 1024:
                    continue
                try:
                    inner = z.read(info)
                except Exception:
                    continue
                text = extract_text(inner, info.filename, "", _zip_depth=depth + 1)
                if (text or "").strip():
                    out.append(f"### {info.filename}\n{text}")
    except Exception:
        return ""
    return "\n\n".join(out).strip()


def _looks_text(raw: bytes) -> bool:
    """Heurística: ¿los bytes son texto legible (no binario)?"""
    sample = raw[:4096]
    if b"\x00" in sample:
        return False
    try:
        sample.decode("utf-8")
        return True
    except UnicodeDecodeError:
        return False


def extract_text(raw: bytes, filename: str = "", mime: str = "", _zip_depth: int = 0) -> str:
    """Devuelve el texto de un archivo subido según su tipo. Nunca lanza.

    Soporta PDF, DOCX, XLSX, PPTX, ZIP (recursivo) e imágenes (OCR opcional).
    Para binarios NUNCA decodifica bytes crudos (evita basura en el RAG): si no
    puede extraer texto, devuelve cadena vacía y el archivo queda solo catalogado.
    """
    name = (filename or "").lower()
    mime = (mime or "").lower()

    if name.endswith(".pdf") or "application/pdf" in mime:
        return _pdf_text(raw)
    if name.endswith(".docx") or "wordprocessingml" in mime:
        return _docx_text(raw)
    if name.endswith(".xlsx") or "spreadsheetml" in mime:
        return _xlsx_text(raw)
    if name.endswith(".pptx") or "presentationml" in mime:
        return _pptx_text(raw)
    if name.endswith(".zip") or "application/zip" in mime:
        return _zip_text(raw, depth=_zip_depth)
    if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif")) or mime.startswith("image/"):
        return _image_ocr(raw)

    # Texto plano, CSV, Markdown, JSON, código… solo si parece texto.
    if _looks_text(raw):
        return raw.decode("utf-8", errors="ignore")
    # Binario desconocido: se cataloga el archivo pero sin texto (sin basura al RAG).
    return ""

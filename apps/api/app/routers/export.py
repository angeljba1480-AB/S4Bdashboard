"""Export endpoints (blueprint backlog: "Export reports" — DOCX/PDF/Markdown).

Generates a professionally formatted Word document (.docx), a PDF, or Markdown
from a conversation transcript or arbitrary report content (proposals, SOWs,
cyber diagnostics, executive summaries). The Word renderer applies the tenant's
white-label branding (name, color, tagline) with a cover header, styled
sections, a data table, and a confidential footer with page numbers.
"""
from __future__ import annotations

import io
import re
from datetime import datetime

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response, StreamingResponse
from pydantic import BaseModel
from sqlmodel import Session, select

from ..auth import get_current_tenant, get_current_user
from ..db import get_session
from ..models import Conversation, Message, Tenant, User

router = APIRouter(prefix="/export", tags=["export"])

DEFAULT_BRAND = "MaestroAI"
_MUTED = (0x6B, 0x72, 0x80)
_INK = (0x11, 0x18, 0x27)
DOCX_MEDIA = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MEDIA = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MEDIA = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


class ReportRequest(BaseModel):
    title: str
    content: str
    format: str = "pdf"  # pdf | md | docx


# --------------------------------------------------------------------------- #
# Branding helpers
# --------------------------------------------------------------------------- #
def _brand_of(tenant: Tenant | None) -> dict:
    return {
        "name": (getattr(tenant, "brand_name", None) or DEFAULT_BRAND),
        "color": (getattr(tenant, "brand_color", None) or "#7C3AED"),
        "tagline": (getattr(tenant, "brand_tagline", None) or ""),
    }


def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = (h or "#7C3AED").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return (0x7C, 0x3A, 0xED)


# --------------------------------------------------------------------------- #
# Markdown-lite parsing (shared shape: list of (heading, body) blocks)
# --------------------------------------------------------------------------- #
def _strip_md(text: str) -> str:
    return text.replace("**", "").replace("`", "")


# --------------------------------------------------------------------------- #
# Word (.docx) — the professional deliverable
# --------------------------------------------------------------------------- #
def _render_docx(
    title: str,
    blocks: list[tuple[str, str]],
    *,
    brand: str = DEFAULT_BRAND,
    brand_color: str = "#7C3AED",
    tagline: str = "",
    meta: dict | None = None,
) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt, RGBColor

    accent = RGBColor(*_hex_to_rgb(brand_color))
    muted = RGBColor(*_MUTED)
    ink = RGBColor(*_INK)

    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    def _heading(text: str, size: int = 14, color: RGBColor = accent, space_before: int = 12):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(space_before)
        p.paragraph_format.space_after = Pt(4)
        r = p.add_run(_strip_md(text).strip())
        r.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = color
        return p

    def _inline(p, text: str, muted_italic: bool = False):
        t = text.strip()
        if len(t) >= 2 and t.startswith("_") and t.endswith("_"):
            text, muted_italic = t[1:-1], True
        for part in re.split(r"(\*\*[^*]+\*\*)", text):
            if not part:
                continue
            run = p.add_run(part[2:-2] if part.startswith("**") and part.endswith("**") else part)
            if part.startswith("**") and part.endswith("**"):
                run.bold = True
            if muted_italic:
                run.italic = True
                run.font.color.rgb = muted

    def _body(text: str):
        for raw in (text or "").split("\n"):
            line = raw.rstrip()
            if not line.strip():
                continue
            if line.startswith("### "):
                _heading(line[4:], 12, ink, 8)
            elif line.startswith("## "):
                _heading(line[3:], 13, ink, 8)
            elif line.startswith("# "):
                _heading(line[2:], 14, ink, 8)
            elif line.lstrip().startswith(("- ", "* ")):
                _inline(doc.add_paragraph(style="List Bullet"), line.lstrip()[2:])
            elif re.match(r"^\s*\d+\.\s+", line):
                _inline(doc.add_paragraph(style="List Number"), re.sub(r"^\s*\d+\.\s+", "", line))
            else:
                _inline(doc.add_paragraph(), line)

    # --- Cover header --------------------------------------------------------
    bp = doc.add_paragraph()
    br = bp.add_run(brand.upper())
    br.bold = True
    br.font.size = Pt(12)
    br.font.color.rgb = accent
    if tagline:
        tr = bp.add_run("   ·   " + tagline)
        tr.italic = True
        tr.font.size = Pt(9)
        tr.font.color.rgb = muted

    tp = doc.add_paragraph()
    tp.paragraph_format.space_before = Pt(2)
    tp.paragraph_format.space_after = Pt(2)
    tr = tp.add_run(title)
    tr.bold = True
    tr.font.size = Pt(22)
    tr.font.color.rgb = ink

    # Accent rule under the title
    pPr = tp._p.get_or_add_pPr()
    pbdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "12")
    bottom.set(qn("w:space"), "4")
    bottom.set(qn("w:color"), brand_color.lstrip("#"))
    pbdr.append(bottom)
    pPr.append(pbdr)

    metabits = [f"Generado {datetime.utcnow():%d/%m/%Y %H:%M} UTC"]
    for k, v in (meta or {}).items():
        if v:
            metabits.append(f"{k}: {v}")
    mp = doc.add_paragraph()
    mr = mp.add_run("   ·   ".join(metabits))
    mr.italic = True
    mr.font.size = Pt(9)
    mr.font.color.rgb = muted
    doc.add_paragraph()

    # --- Sections ------------------------------------------------------------
    for heading, body in blocks:
        if heading and heading.lower().startswith("datos proporcionados"):
            _heading(heading)
            rows = [ln.split(":", 1) for ln in (body or "").splitlines() if ":" in ln]
            if rows:
                table = doc.add_table(rows=0, cols=2)
                try:
                    table.style = "Light List Accent 1"
                except Exception:
                    table.style = "Table Grid"
                for k, v in rows:
                    cells = table.add_row().cells
                    kr = cells[0].paragraphs[0].add_run(k.strip().capitalize())
                    kr.bold = True
                    cells[1].paragraphs[0].add_run(v.strip())
            doc.add_paragraph()
            continue
        if heading:
            _heading(heading)
        _body(body)

    # --- Footer (brand · confidential · page number) -------------------------
    footer = doc.sections[0].footer
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = fp.add_run(f"{brand}  ·  Documento confidencial  ·  Página ")
    fr.font.size = Pt(8)
    fr.font.color.rgb = muted
    run = fp.add_run()
    run.font.size = Pt(8)
    run.font.color.rgb = muted
    for kind, txt in (("begin", None), ("instr", "PAGE"), ("end", None)):
        el = OxmlElement("w:fldChar") if kind != "instr" else OxmlElement("w:instrText")
        if kind == "instr":
            el.set(qn("xml:space"), "preserve")
            el.text = txt
        else:
            el.set(qn("w:fldCharType"), kind)
        run._r.append(el)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# PDF / Markdown (kept simple; branding fixed)
# --------------------------------------------------------------------------- #
def _render_pdf(title: str, blocks: list[tuple[str, str]], brand: str = DEFAULT_BRAND) -> bytes:
    """blocks = list of (heading, body). Pure-Python via reportlab."""
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER, title=title)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]),
             Paragraph(f"Generado {datetime.utcnow():%Y-%m-%d %H:%M} UTC · {brand}",
                       styles["Italic"]),
             Spacer(1, 16)]
    for heading, body in blocks:
        if heading:
            story.append(Paragraph(heading, styles["Heading3"]))
        for line in (body or "").split("\n"):
            story.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;") or "&nbsp;", styles["BodyText"]))
        story.append(Spacer(1, 10))
    doc.build(story)
    return buf.getvalue()


def _render_pptx(
    title: str, blocks: list[tuple[str, str]], *,
    brand: str = DEFAULT_BRAND, brand_color: str = "#7C3AED", tagline: str = "",
) -> bytes:
    """One title slide + one content slide per block (heading → bullets)."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    accent = RGBColor(*_hex_to_rgb(brand_color))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    tb = cover.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5)).text_frame
    tb.word_wrap = True
    run = tb.paragraphs[0].add_run(); run.text = brand.upper()
    run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = accent
    p = tb.add_paragraph(); r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True
    sub = tb.add_paragraph(); sr = sub.add_run()
    sr.text = (tagline + "  ·  " if tagline else "") + f"Generado {datetime.utcnow():%d/%m/%Y} UTC"
    sr.font.size = Pt(12); sr.font.color.rgb = RGBColor(*_MUTED)

    for heading, body in blocks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        head = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1)).text_frame
        hr = head.paragraphs[0].add_run(); hr.text = _strip_md(heading or title)
        hr.font.size = Pt(26); hr.font.bold = True; hr.font.color.rgb = accent
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.4)).text_frame
        box.word_wrap = True
        first = True
        for raw in (body or "").split("\n"):
            line = _strip_md(raw).strip().lstrip("-*").strip()
            if not line:
                continue
            para = box.paragraphs[0] if first else box.add_paragraph()
            first = False
            run = para.add_run(); run.text = line; run.font.size = Pt(16)
            para.level = 0
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_xlsx(title: str, blocks: list[tuple[str, str]], brand: str = DEFAULT_BRAND) -> bytes:
    """A single sheet: section headers + body lines; key:value rows become 2 columns."""
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "Reporte"
    ws.column_dimensions["A"].width = 32
    ws.column_dimensions["B"].width = 80

    head_fill = PatternFill("solid", fgColor="0A1F44")
    ws.append([title]); ws["A1"].font = Font(bold=True, size=14)
    ws.append([f"Generado {datetime.utcnow():%Y-%m-%d %H:%M} UTC · {brand}"])
    ws["A2"].font = Font(italic=True, size=9, color="6B7280")
    ws.append([])

    for heading, body in blocks:
        if heading:
            r = ws.max_row + 1
            ws.append([_strip_md(heading)])
            c = ws.cell(row=r, column=1)
            c.font = Font(bold=True, color="FFFFFF"); c.fill = head_fill
        for raw in (body or "").split("\n"):
            line = raw.rstrip()
            if not line.strip():
                continue
            if ":" in line and len(line.split(":", 1)[0]) < 40:
                k, v = line.split(":", 1)
                ws.append([_strip_md(k).strip(), _strip_md(v).strip()])
            else:
                ws.append(["", _strip_md(line).strip()])
                ws.cell(row=ws.max_row, column=2).alignment = Alignment(wrap_text=True)
        ws.append([])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _render_md(title: str, blocks: list[tuple[str, str]], brand: str = DEFAULT_BRAND) -> str:
    out = [f"# {title}", f"_Generado {datetime.utcnow():%Y-%m-%d %H:%M} UTC · {brand}_", ""]
    for heading, body in blocks:
        if heading:
            out.append(f"## {heading}")
        out.append(body or "")
        out.append("")
    return "\n".join(out)


# --------------------------------------------------------------------------- #
# Shared dispatcher
# --------------------------------------------------------------------------- #
def deliver(title: str, blocks: list[tuple[str, str]], fmt: str, fname: str,
            tenant: Tenant | None = None, meta: dict | None = None):
    b = _brand_of(tenant)
    if fmt == "md":
        return Response(_render_md(title, blocks, brand=b["name"]), media_type="text/markdown",
                        headers={"Content-Disposition": f'attachment; filename="{fname}.md"'})
    if fmt == "docx":
        data = _render_docx(title, blocks, brand=b["name"], brand_color=b["color"],
                            tagline=b["tagline"], meta=meta)
        return StreamingResponse(io.BytesIO(data), media_type=DOCX_MEDIA,
                                 headers={"Content-Disposition": f'attachment; filename="{fname}.docx"'})
    if fmt == "pptx":
        data = _render_pptx(title, blocks, brand=b["name"], brand_color=b["color"], tagline=b["tagline"])
        return StreamingResponse(io.BytesIO(data), media_type=PPTX_MEDIA,
                                 headers={"Content-Disposition": f'attachment; filename="{fname}.pptx"'})
    if fmt == "xlsx":
        data = _render_xlsx(title, blocks, brand=b["name"])
        return StreamingResponse(io.BytesIO(data), media_type=XLSX_MEDIA,
                                 headers={"Content-Disposition": f'attachment; filename="{fname}.xlsx"'})
    pdf = _render_pdf(title, blocks, brand=b["name"])
    return StreamingResponse(io.BytesIO(pdf), media_type="application/pdf",
                             headers={"Content-Disposition": f'attachment; filename="{fname}.pdf"'})


@router.post("/report")
def export_report(
    body: ReportRequest,
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant),
):
    return deliver(body.title, [("", body.content)], body.format, body.title,
                   tenant=tenant, meta={"Preparado por": user.name})


class ToCloudRequest(BaseModel):
    title: str
    content: str
    provider: str = "google"   # google | microsoft


@router.post("/to-cloud")
def export_to_cloud(
    body: ToCloudRequest,
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant),
    session: Session = Depends(get_session),
) -> dict:
    """Guarda el reporte en la nube del cliente: Google Docs (gdocs.create) u
    OneDrive (onedrive.upload), usando su cuenta conectada (toolkit Google/MS)."""
    from ..integrations import actions_exec, token_store
    prov = "microsoft" if body.provider == "microsoft" else "google"
    # Solo la cuenta del PROPIO usuario (no escribir en el Drive de otro del tenant).
    tok = token_store.get_valid_access_token(session, tenant, user.id, prov)
    if not tok:
        raise HTTPException(status_code=400,
                            detail=f"Conecta TU cuenta {prov} en Integraciones para guardar ahí.")
    name = (body.title or "Reporte MaestroAI").strip()
    action = "gdocs.create" if prov == "google" else "onedrive.upload"
    params = {"title": name, "name": f"{name}.txt", "content": body.content}
    try:
        detail = actions_exec.execute(action, tok, params)
    except Exception as exc:  # pragma: no cover - red/credenciales
        raise HTTPException(status_code=400, detail=f"No se pudo crear en {prov}: {exc}")
    return {"ok": True, "provider": prov, "detail": detail}


@router.get("/conversation/{conversation_id}")
def export_conversation(
    conversation_id: str,
    format: str = "pdf",
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant),
    session: Session = Depends(get_session),
):
    conv = session.get(Conversation, conversation_id)
    if not conv or conv.tenant_id != tenant.id:
        raise HTTPException(status_code=404, detail="Conversación no encontrada")
    msgs = session.exec(
        select(Message).where(Message.conversation_id == conv.id).order_by(Message.created_at)
    ).all()
    blocks = [
        (f"{'Usuario' if m.role == 'user' else 'Asistente'}"
         f"{f' · {m.model_used} ({m.route.value})' if m.role == 'assistant' else ''}",
         m.content_redacted)
        for m in msgs
    ]
    title = conv.title or "Conversación"
    return deliver(title, blocks, format, conv.id, tenant=tenant,
                   meta={"Preparado por": user.name})

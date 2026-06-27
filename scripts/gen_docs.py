"""Documentación EXTENSA y profunda de MaestroAI → Word (.docx) y PowerPoint (.pptx).

Uso:  OUT_DIR=docs/generados python scripts/gen_docs.py
Estructura: por cada capacidad → Capacidad · Logro/valor · Caso de uso · Pasos a
seguir. Incluye guías técnicas de configuración paso a paso y apéndices.
"""
from __future__ import annotations

import os

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PColor

OUT = os.environ.get("OUT_DIR", ".")
VIOLET = RGBColor(0x6D, 0x28, 0xD9); SLATE = RGBColor(0x33, 0x41, 0x55)
GREY = RGBColor(0x64, 0x74, 0x8B)
PV = PColor(0x6D, 0x28, 0xD9); PW = PColor(0xFF, 0xFF, 0xFF)
PS = PColor(0x33, 0x41, 0x55); PG = PColor(0x64, 0x74, 0x8B)

# ============================================================ CAPACIDADES (deep)
CAPS = [
    {
        "t": "Enrutador de Privacidad (Privacy Model Router)",
        "cap": "Componente que, por CADA petición, clasifica la sensibilidad del dato "
               "(Público/Interno/Confidencial/Restringido), detecta PII (RFC, CURP, CLABE, "
               "tarjetas, correos, salud, secretos), minimiza y redacta el contexto, y decide "
               "automáticamente la ruta del modelo: local, VPC, abierta (NaN), premium o BLOQUEO.",
        "logro": "Permite usar IA potente sin exponer información sensible ni depender del criterio "
                 "del usuario. Es el control que convierte «IA» en «IA gobernada»: soberanía de datos, "
                 "cumplimiento (LFPDPPP) y trazabilidad por diseño.",
        "caso": "Un analista pega un contrato con cláusula NDA y datos bancarios. El router lo clasifica "
                "como Confidencial con PII, lo mantiene en la ruta privada (VPC/local), redacta los "
                "identificadores y deja constancia auditable — sin que el analista decida nada.",
        "pasos": [
            "El usuario solo escribe su consulta o sube su documento; no elige modelo.",
            "El sistema clasifica y muestra un banner con la clasificación y la ruta real usada.",
            "Si el dato es Restringido, se procesa local y nunca sale; si es Confidencial, VPC o local.",
            "Para revisar decisiones: Auditoría muestra ruta, modelo, tokens, costo y razón por evento.",
        ],
    },
    {
        "t": "RAG documental por área",
        "cap": "Repositorio de documentos por área y categoría con tratamiento auto-detectado "
               "(público/interno/confidencial/restringido) y editable. Trocea (chunking con solape), "
               "calcula embeddings y cifra los fragmentos; la recuperación se acota por tenant, "
               "categoría y áreas visibles del usuario, y responde citando la fuente.",
        "logro": "Convierte el conocimiento disperso (propuestas, licitaciones, ISO, manuales) en respuestas "
                 "confiables con cita y respetando permisos. Reduce el «inventar» del modelo y el riesgo de "
                 "fuga entre áreas.",
        "caso": "El área de Ventas pregunta «¿qué incluimos en la última propuesta a ACME?» y obtiene la "
                "respuesta citando el documento exacto, sin ver documentos de RH o Finanzas.",
        "pasos": [
            "Documentos → Subir: asigna área y categoría; revisa el tratamiento detectado.",
            "El contenido se cifra e indexa automáticamente en el RAG.",
            "En Chat con fuentes elige el modo de contexto (todo / elegir documentos / sin contexto).",
            "La respuesta cita el documento y el fragmento; los permisos por área se aplican siempre.",
        ],
    },
    {
        "t": "Reranking del RAG (precisión)",
        "cap": "Tras recuperar candidatos por embeddings, los reordena por relevancia real query↔documento "
               "con un cross-encoder (NaN Qwen3-Reranker, endpoint /rerank). Patrón embedding → rerank → LLM.",
        "logro": "Sube notablemente la precisión de las respuestas con fuentes: el modelo recibe primero los "
                 "fragmentos verdaderamente relevantes, no solo los «parecidos» por vector.",
        "caso": "En una base de 10.000 fragmentos, la pregunta recupera 20 candidatos y el reranker deja los "
                "4 más pertinentes arriba, mejorando la calidad de la cita y la respuesta.",
        "pasos": [
            "Admin → Eficiencia de tokens → activa «Reranking RAG».",
            "Requiere el proveedor abierto (NaN) configurado (usa su Base URL + API key).",
            "Si el reranker falla, el sistema cae limpio al orden por coseno (sin romperse).",
        ],
    },
    {
        "t": "Cascada NaN → premium y eficiencia de tokens",
        "cap": "Empieza SIEMPRE con el modelo abierto (NaN), barato y rápido. Escala a premium solo a demanda "
               "(«máxima precisión») o cuando la respuesta barata es insuficiente. Antes de pagar premium, "
               "condensa el contexto grande con el modelo barato y aplica un tope de gasto por consulta.",
        "logro": "Calidad cuando importa, costo bajo por defecto. Evita pagar premium en todo y reduce los "
                 "tokens enviados (condensación), con ahorro acumulado visible.",
        "caso": "Un PDF de 40 páginas se condensa con NaN a un extracto; si el usuario pide «máxima precisión», "
                "premium recibe solo ese extracto y refina — pagando una fracción.",
        "pasos": [
            "No requiere acción del usuario para el flujo base (todo empieza en NaN).",
            "Activa «Máxima precisión» en el chat para forzar el refinamiento premium.",
            "Admin → Eficiencia: define umbral de condensación y tope de tokens por consulta.",
            "Si no hay premium configurado, todo corre en NaN; sin NaN, modo demostración (mock).",
        ],
    },
    {
        "t": "Casos de uso (generación de entregables)",
        "cap": "Plantillas que generan documentos profesionales (propuesta comercial, carta, reporte, "
               "licitación, reporte por industria) con un paso universal objetivo/notas/formato y grounding "
               "por categoría del RAG. Exporta a Word, PDF, Markdown, PPTX y XLSX.",
        "logro": "Pasa de horas a minutos en entregables repetitivos, con marca de la empresa, estructura "
                 "consistente y datos citados. «Elige qué lograr, da lo mínimo; tú solo apruebas.»",
        "caso": "Un ejecutivo elige «Propuesta comercial», pone cliente/servicio/monto, revisa el borrador y "
                "descarga un .docx con portada, secciones y pie confidencial.",
        "pasos": [
            "Casos de uso → elige la receta.",
            "Llena los datos mínimos (objetivo, notas, formato de salida).",
            "Revisa el borrador (vía router de privacidad) y pulsa «Aprobar y generar».",
            "Descarga en Word/PDF/Markdown/PPTX/XLSX.",
        ],
    },
    {
        "t": "Agentes verticales",
        "cap": "Agentes especializados por dominio (Document Intelligence, Cyber Diagnostic, Proposal/SOW, "
               "Executive Copilot) con política de privacidad propia; cada respuesta pasa por el router y cita.",
        "logro": "Experiencia experta lista para usar por área, sin construir prompts desde cero, con la misma "
                 "gobernanza de privacidad.",
        "caso": "El equipo de seguridad usa Cyber Diagnostic para un diagnóstico y roadmap citando las políticas "
                "internas.",
        "pasos": [
            "Agentes → elige el agente y conversa.",
            "Cada respuesta cita sus fuentes y respeta el ruteo de privacidad.",
            "El Admin puede crear agentes personalizados por área.",
        ],
    },
    {
        "t": "Chat con fuentes",
        "cap": "Chat con 3 modos de contexto (sin contexto / todo el RAG / elegir documentos), toggles de "
               "«máxima precisión» (cascada) y «usar memoria», banner con clasificación y ruta real, exportable.",
        "logro": "Una sola interfaz para preguntar con o sin documentos, con control de precisión y costo, y "
                 "trazabilidad de cada respuesta.",
        "caso": "Un usuario activa «elegir documentos», selecciona 3 archivos y pregunta; la respuesta se basa "
                "solo en esos, con cita.",
        "pasos": [
            "Chat con fuentes → elige el modo de contexto.",
            "Opcional: activa «Máxima precisión» y/o «Usar memoria».",
            "Envía; revisa la ruta en el banner y exporta a PDF/Markdown si lo necesitas.",
        ],
    },
    {
        "t": "Notebooks (estilo NotebookLM privado)",
        "cap": "Espacios con fuentes acotadas donde se pregunta con cita y se generan artefactos (resumen, "
               "FAQ, guía, briefing, cronología) solo sobre esos documentos.",
        "logro": "Estudio profundo de un conjunto de documentos sin que el modelo se «desvíe» a otros, con "
                 "salidas reutilizables.",
        "caso": "Para una licitación se crea un notebook con las bases y se genera una FAQ y una cronología.",
        "pasos": [
            "Notebooks → crea uno y añade fuentes del RAG.",
            "Haz preguntas (responde citando) o genera un artefacto.",
        ],
    },
    {
        "t": "Memoria y etiquetas",
        "cap": "Memoria persistente de trabajos con búsqueda semántica y por tags; auto-captura al completar "
               "casos; recall en el chat respetando permisos por área.",
        "logro": "El sistema «recuerda» trabajos previos («¿recuerdas el trabajo C?») y los reutiliza, evitando "
                 "rehacer y perdiendo menos contexto entre sesiones.",
        "caso": "Semanas después, el usuario pide continuar «la propuesta C» y el sistema recupera lo hecho.",
        "pasos": [
            "Guarda resultados con «Guardar en memoria» (o se auto-capturan al completar casos).",
            "Organiza con etiquetas y busca por texto o tag en Memoria.",
            "Activa «Usar memoria» en el chat para responder con base en trabajos previos.",
        ],
    },
    {
        "t": "Generar imágenes (texto → imagen)",
        "cap": "Generación de imágenes por ruta estándar OpenAI (/images/generations) con relación de aspecto "
               "y variantes; galería por área; el prompt se redacta de PII; todo auditado.",
        "logro": "Material visual gobernado dentro de la plataforma, con copia propia y permisos por área.",
        "caso": "Marketing genera variantes de una imagen para una campaña; quedan en la galería del área.",
        "pasos": [
            "Generar imágenes → escribe el prompt, elige aspecto y variantes.",
            "Requiere un proveedor que exponga /images/generations (la API de NaN no lo expone hoy).",
        ],
    },
    {
        "t": "Toolkit de acciones (Google / Microsoft)",
        "cap": "Ejecuta tareas reales en las herramientas del usuario vía su OAuth: enviar correo, crear "
               "eventos, append a Sheets/Excel, publicar en Teams, buscar en SharePoint, listar OneDrive, "
               "leer Sheets/Excel/Calendar. Lecturas inmediatas; escrituras con aprobación humana.",
        "logro": "La IA no solo redacta: actúa, con un control de aprobación («tú apruebas») y «Permitir "
                 "siempre» revocable, todo auditado.",
        "caso": "Tras redactar un correo, el usuario pulsa enviar; queda como solicitud pendiente y, al aprobar, "
                "se manda desde su cuenta Outlook.",
        "pasos": [
            "Configura scopes de escritura y reconecta (ver guía de configuración A/B).",
            "Acciones → elige la acción y llena los datos.",
            "Aprueba la escritura (o usa «Permitir siempre»).",
        ],
    },
    {
        "t": "Automatizaciones y n8n",
        "cap": "Conecta un disparador (manual / programado / por evento) con una acción (workflow n8n, caso de "
               "uso, conector o notificar). n8n gestionado se auto-provisiona por tenant (6 workflows base).",
        "logro": "Procesos sin intervención: resúmenes diarios, cobranza semanal, alertas de documento sensible, "
                 "reportes — orquestando los sistemas de la empresa.",
        "caso": "Cada vez que se sube un documento confidencial, una automatización notifica a seguridad.",
        "pasos": [
            "Automatizaciones → Nueva (o usa una plantilla).",
            "Elige disparador y acción («workflow» para n8n con su referencia).",
            "Activa/desactiva o «Ejecutar ahora»; el resultado queda en Auditoría.",
        ],
    },
    {
        "t": "Conectores, webhooks y API pública",
        "cap": "Conectores REST de salida (CRM/ERP/Delivery) con detalle y «ojito» de secretos (auditado); "
               "webhooks entrantes firmados (HMAC-SHA256); API pública /v1 con X-API-Key.",
        "logro": "Integración bidireccional gobernada con el ecosistema de la empresa, sin exponer secretos.",
        "caso": "Un lead capturado en la plataforma se envía a HubSpot vía conector; un sistema externo notifica "
                "eventos por webhook firmado.",
        "pasos": [
            "Integraciones → Conectores: plantilla, endpoint y token (cifrado); Probar.",
            "Webhooks entrantes → crea uno y firma el cuerpo con el secreto.",
            "API pública: genera una API key en Admin y llama /v1 con X-API-Key.",
        ],
    },
    {
        "t": "Fuentes de datos legadas (BD de solo lectura / CSV)",
        "cap": "Importa al RAG desde una base de datos de solo lectura (DSN + SELECT, con denylist DML/CTE y "
               "transacción read-only) o desde un CSV pegado (encabezados + delimitador). Solo ADMIN/DEVOPS.",
        "logro": "Trae el conocimiento de sistemas sin API a la IA gobernada, sin exponer escritura ni romper "
                 "el sistema legado.",
        "caso": "Un ERP antiguo exporta clientes en CSV; se importan al RAG y el chat ya responde sobre ellos.",
        "pasos": [
            "Integraciones → Fuentes de datos: DSN + SELECT → Probar → Importar.",
            "O «Importar CSV»: pega el contenido y el delimitador.",
            "Usa el «ojito» para revelar el DSN (auditado).",
        ],
    },
    {
        "t": "Tableros, Trámites y App Studio",
        "cap": "Tableros a la medida (describe qué medir → widgets en vivo + KPIs manuales); catálogo de "
               "Trámites por país/estado/empresa (MCP) que aterriza respuestas; App Studio para mini-apps.",
        "logro": "Medición y contexto público/privado curado, y la posibilidad de publicar mini-aplicaciones.",
        "caso": "Un gerente arma un tablero de «costo por ruta de modelo» y casos por receta; convierte un "
                "trámite estatal en una propuesta de caso de uso.",
        "pasos": [
            "Tableros → describe lo que quieres medir y agrega widgets.",
            "Trámites y casos → filtra y convierte en propuesta.",
        ],
    },
    {
        "t": "Auditoría, gobernanza y multi-tenant",
        "cap": "Auditoría navegable de cada evento (usuario, ruta, modelo, tokens, costo, sensibilidad, razón) "
               "con export a SIEM (JSONL); permisos jerárquicos por área y licencia; super admin multi-tenant.",
        "logro": "Evidencia y control para cumplimiento y seguridad; cada empresa aislada; cada usuario ve lo suyo.",
        "caso": "Compliance revisa en Auditoría todas las salidas externas del mes y las exporta a su SIEM.",
        "pasos": [
            "Auditoría → filtra por tipo, riesgo, ruta o usuario; abre el detalle.",
            "Export SIEM → descarga JSONL.",
            "Admin → Usuarios: asigna rol y área a cada persona.",
        ],
    },
]

# ============================================================ GUÍAS DE CONFIG
GUIDES = [
    {
        "t": "A. Microsoft 365 / Outlook (OAuth + acciones)",
        "intro": "Habilita el resumen de correo/agenda y el toolkit de acciones de Microsoft. Requiere "
                 "registrar una app en Azure (Entra ID).",
        "pasos": [
            "portal.azure.com → Microsoft Entra ID → App registrations → New registration. Nombre: MaestroAI.",
            "Supported account types: «Accounts in any organizational directory and personal Microsoft accounts».",
            "Redirect URI (Web): https://<tu-api>/oauth/microsoft/callback (la URL real del backend en Render).",
            "API permissions → Microsoft Graph → Delegated → lectura: User.Read, Mail.Read, Calendars.Read, offline_access.",
            "Acciones (escritura): Mail.Send, Calendars.ReadWrite, ChannelMessage.Send, Files.ReadWrite.All, Sites.Read.All.",
            "(Empresa) Grant admin consent.",
            "Certificates & secrets → New client secret → copia el Value (solo se ve una vez) y el Application (client) ID.",
            "Configura las variables MICROSOFT_* en Render → Save, rebuild & deploy.",
            "Cada usuario Reconecta una vez (Integraciones → Conectar correo) para otorgar los nuevos permisos.",
        ],
        "env": [
            "MICROSOFT_OAUTH_ENABLED = true",
            "MICROSOFT_CLIENT_ID     = <Application (client) ID>",
            "MICROSOFT_CLIENT_SECRET = <el Value del secreto>",
            "MICROSOFT_TENANT        = common",
            "MICROSOFT_REDIRECT_URI  = https://<tu-api>/oauth/microsoft/callback",
            "APP_PUBLIC_URL          = https://plataforma.maestroai.mx",
        ],
        "nota": "La REDIRECT_URI en Render debe coincidir EXACTAMENTE con la registrada en Azure. Los tokens se "
                "guardan cifrados (AES-256-GCM) y se refrescan solos.",
    },
    {
        "t": "B. Google Workspace / Gmail (OAuth + acciones)",
        "intro": "Habilita Gmail/Calendar/Drive/Sheets para resumen, contexto y acciones.",
        "pasos": [
            "console.cloud.google.com → crea proyecto → APIs & Services.",
            "Habilita: Gmail API, Google Calendar API, Google Drive API, Google Sheets API.",
            "OAuth consent screen: External; agrega tu correo como test user mientras esté en pruebas.",
            "Scopes: gmail.readonly, calendar.events, drive.readonly y (acciones) gmail.send, spreadsheets.",
            "Credentials → Create credentials → OAuth client ID → Web application.",
            "Authorized redirect URI: https://<tu-api>/oauth/google/callback.",
            "Configura GOOGLE_* en Render y reconecta desde el portal.",
        ],
        "env": [
            "GOOGLE_OAUTH_ENABLED = true",
            "GOOGLE_CLIENT_ID     = <client id>.apps.googleusercontent.com",
            "GOOGLE_CLIENT_SECRET = <client secret>",
            "GOOGLE_REDIRECT_URI  = https://<tu-api>/oauth/google/callback",
        ],
        "nota": "Gmail/Outlook ya no aceptan contraseña normal por IMAP; por eso usan OAuth.",
    },
    {
        "t": "C. Cualquier otro correo — IMAP",
        "intro": "Yahoo, iCloud, Zoho, hosting o correo de empresa. No requiere registrar ninguna app.",
        "pasos": [
            "Integraciones → Conectar correo → «Otro correo (IMAP)».",
            "Elige el proveedor (preset) o escribe host/puerto.",
            "Escribe el correo + contraseña (muchos exigen «contraseña de aplicación» con 2FA activado).",
            "Queda conectado; la contraseña se guarda cifrada y se usa solo para leer la bandeja.",
        ],
        "env": [],
        "nota": "IMAP trae correo (no calendario; eso solo Outlook/Gmail por OAuth).",
    },
    {
        "t": "D. Proveedor abierto — NaN Builders (ruta open)",
        "intro": "Modelo barato/rápido y base de la cascada, el rerank y la condensación. API compatible OpenAI.",
        "pasos": [
            "Obtén tu API key de NaN (nan.builders/docs/getting-started).",
            "UI: Admin → Modelos externos → Abierto → Base URL + modelo + API key → Guardar → Probar conexión.",
            "O configura OPEN_* en Render.",
            "Modelo recomendado: qwen3.6 (principal). Para razonamiento largo: deepseek-v4-flash.",
        ],
        "env": [
            "OPEN_ENABLED  = true",
            "OPEN_API_KEY  = sk-...",
            "OPEN_BASE_URL = https://api.nan.builders/v1",
            "OPEN_MODEL    = qwen3.6",
        ],
        "nota": "Límites NaN: 60 rpm, 3 en paralelo, 1.5M tpm por modelo de chat. Enterprise (Helmcode): "
                "api.helmcode.com/v1.",
    },
    {
        "t": "E. Proveedor premium (OpenAI / Claude / Gemini compatible)",
        "intro": "Refinamiento de máxima precisión en la cascada. Solo se usa como escalada.",
        "pasos": [
            "Admin → Modelos externos → Premium → Base URL + modelo + API key → Guardar.",
            "Pulsa «Probar conexión» para validar (latencia + muestra, o el error).",
            "O configura PREMIUM_* en Render.",
        ],
        "env": [
            "PREMIUM_ENABLED  = true",
            "PREMIUM_API_KEY  = <api key>",
            "PREMIUM_BASE_URL = https://api.openai.com/v1   # o el endpoint compatible de tu proveedor",
            "PREMIUM_MODEL    = gpt-4o   # o claude-..., gemini-...",
        ],
        "nota": "Si no se configura, la cascada se queda en NaN (no se promete premium).",
    },
    {
        "t": "F. Ollama — ruta local privada (datos confidenciales)",
        "intro": "Procesa lo Confidencial/Restringido en tu propia máquina/servidor; nunca sale a la nube.",
        "pasos": [
            "Instala Ollama (ollama.com/download o brew install ollama) y arráncalo (ollama serve).",
            "Descarga un modelo: ollama pull llama3.1 (o qwen2.5:3b para CPU). El nombre exacto = LOCAL_MODEL.",
            "Exponlo con URL pública (prueba: cloudflared tunnel --url http://localhost:11434; prod: servidor HTTPS).",
            "Configura LOCAL_* en Render con la URL del túnel/servidor + /v1.",
            "Mientras no haya local, ALLOW_CLOUD_FALLBACK=true hace que lo sensible use NaN real en vez del simulador.",
        ],
        "env": [
            "LOCAL_ENABLED  = true",
            "LOCAL_BASE_URL = https://<tu-tunel-o-servidor>/v1",
            "LOCAL_MODEL    = llama3.1",
        ],
        "nota": "Al conectar Ollama, lo sensible vuelve a la ruta local automáticamente (puedes apagar el fallback).",
    },
    {
        "t": "G. n8n (workflows / automatización)",
        "intro": "Gestionado por defecto (cero config). BYO opcional para usar tu propio n8n.",
        "pasos": [
            "Admin → n8n: Webhook Base URL + API key → Guardar (verás el motor y el origen).",
            "O configura N8N_* en Render. N8N_AUTO_PROVISION crea los workflows por tenant.",
            "En Automatizaciones usa acción «workflow» con el nombre del workflow como referencia.",
        ],
        "env": [
            "N8N_ENABLED          = true",
            "N8N_WEBHOOK_BASE_URL = https://<tu-n8n>/webhook",
            "N8N_API_BASE_URL     = https://<tu-n8n>/api/v1",
            "N8N_API_KEY          = <api key de n8n>",
            "N8N_AUTO_PROVISION   = true",
        ],
        "nota": "Regla de integración: webhooks para eventos, REST (conectores) para comandos, n8n para el resto.",
    },
    {
        "t": "H. Reranking del RAG",
        "intro": "Mejora la precisión reordenando los candidatos recuperados.",
        "pasos": [
            "Configura primero el proveedor abierto (NaN) — guía D.",
            "Admin → Eficiencia de tokens → activa «Reranking RAG».",
            "O configura RERANK_* en Render.",
        ],
        "env": [
            "RERANK_ENABLED    = true",
            "RERANK_MODEL      = rerank",
            "RERANK_CANDIDATES = 20",
        ],
        "nota": "Si el reranker falla, se usa el orden por coseno (sin romperse).",
    },
    {
        "t": "I. Proteger ramas (GitHub) — gobernanza del código",
        "intro": "Evita pushes directos a producción; exige PR + CI verde. Requiere admin del repo.",
        "pasos": [
            "Repo → Settings → Branches → Add branch ruleset (o Add branch protection rule).",
            "Aplica a main y qa.",
            "Activa: Require a pull request before merging.",
            "Activa: Require status checks to pass → selecciona «API · pytest» y «Web · build».",
            "(Opcional) Require branches to be up to date before merging. Guarda.",
        ],
        "env": [],
        "nota": "Con esto, nadie mete código a qa/main sin PR + CI en verde — el flujo dev → qa → main.",
    },
    {
        "t": "J. Núcleo (Render) + Portal (Vercel) + Base de datos",
        "intro": "Variables base del backend y del portal, y la base de datos persistente.",
        "pasos": [
            "Render → Environment: SECRET_KEY (32+), DATABASE_URL (Supabase Session Pooler), CORS_ORIGINS, APP_ENV=production.",
            "Vercel → NEXT_PUBLIC_API_BASE_URL = URL del backend (sin / final); Root Directory apps/web; Production Branch main.",
            "Supabase: usa la cadena Session Pooler (IPv4); el backend crea tablas y siembra el tenant demo al primer arranque.",
            "Plan Render Starter para que el backend no se «duerma»; salud en /health.",
        ],
        "env": [
            "SECRET_KEY        = <32+ caracteres>",
            "DATABASE_URL      = postgresql://...supabase... (Session Pooler)",
            "CORS_ORIGINS      = https://plataforma.maestroai.mx",
            "APP_ENV           = production",
            "NEXT_PUBLIC_API_BASE_URL = https://<tu-api>   # en Vercel",
        ],
        "nota": "Cada entorno (dev/qa/prod) usa sus propias variables y BD; las llaves de prod no se reutilizan en QA.",
    },
]

NAN_MODELS = [
    ("qwen3.6", "MoE 35B/3B", "256K", "Chat, tools, visión, reasoning — principal"),
    ("deepseek-v4-flash", "MoE 284B/21B", "1M", "Chat, tools, reasoning (reasoning_effort)"),
    ("mimo-v2.5", "MoE 310B/15B", "1M", "Omnimodal (texto+imagen+audio)"),
    ("gemma4", "MoE 26B/4B", "256K", "Chat, visión, reasoning (opt-in)"),
    ("qwen3-embedding", "8B", "—", "Embeddings 4096-dim, 100+ idiomas"),
    ("rerank", "Qwen3-Reranker-8B", "—", "Reranking RAG (/rerank)"),
    ("kokoro", "TTS 82M", "—", "Text-to-speech, voces ES"),
    ("whisper", "large-v3", "—", "Speech-to-text, 99+ idiomas"),
]
ENDPOINTS = {
    "Autenticación / cuenta": ["POST /auth/login", "GET /me", "GET /account"],
    "Chat y casos": ["POST /chat", "POST /chat/preview", "GET/POST /recipes", "POST /recipes/{id}/start", "POST /v1/cases/{recipe_id}/run"],
    "Documentos y RAG": ["GET/POST /documents", "POST /documents/upload", "GET/POST /documents/categories", "GET /drive/files", "POST /drive/import"],
    "Notebooks / Memoria": ["GET/POST /notebooks", "POST /notebooks/{id}/ask", "GET/POST /memory", "GET /memory/tags"],
    "Imágenes": ["GET /images/config", "POST /images/generate", "GET /images", "GET /images/{id}/data"],
    "Acciones": ["GET /actions", "POST /actions/run", "POST /actions/requests/{id}/approve", "GET/DELETE /actions/grants"],
    "Integraciones": ["GET/POST /integrations/connectors", "POST /integrations/connectors/{id}/test", "GET /integrations/connectors/{id}/reveal", "GET/POST /datasources", "POST /datasources/import-csv", "POST /oauth/{provider}/authorize"],
    "Automatizaciones": ["GET/POST /automations", "POST /automations/from-template", "POST /automations/{id}/run", "POST /automations/run-due"],
    "Admin": ["GET/PUT /admin/providers", "POST /admin/providers/{route}/test", "GET/PUT /admin/efficiency", "GET/POST /admin/users", "GET /admin/n8n"],
    "Auditoría / salud": ["GET /audit", "GET /audit/export", "GET /usage", "GET /health", "POST /v1/webhooks/{id}"],
}


# ===================================================================== WORD
def build_docx(path: str) -> None:
    doc = Document()
    base = doc.styles["Normal"]; base.font.name = "Calibri"; base.font.size = Pt(10.5)
    for i, sz in ((1, 18), (2, 14), (3, 12)):
        st = doc.styles[f"Heading {i}"]; st.font.color.rgb = VIOLET; st.font.size = Pt(sz)

    def para(text, bold=False, color=SLATE, size=10.5, italic=False):
        p = doc.add_paragraph(); r = p.add_run(text); r.bold = bold; r.italic = italic
        r.font.color.rgb = color; r.font.size = Pt(size); return p

    def label(lbl, text):
        p = doc.add_paragraph(); r = p.add_run(lbl + ": "); r.bold = True; r.font.color.rgb = VIOLET
        r.font.size = Pt(10.5); rr = p.add_run(text); rr.font.color.rgb = SLATE; rr.font.size = Pt(10.5)

    def steps(items, numbered=True):
        style = "List Number" if numbered else "List Bullet"
        for it in items:
            p = doc.add_paragraph(style=style); p.add_run(it)

    def code(lines):
        for ln in lines:
            p = doc.add_paragraph(); r = p.add_run(ln); r.font.name = "Consolas"; r.font.size = Pt(9)
            r.font.color.rgb = SLATE

    def table(headers, rows):
        t = doc.add_table(rows=1, cols=len(headers)); t.alignment = WD_TABLE_ALIGNMENT.CENTER
        t.style = "Light Grid Accent 1"
        for i, h in enumerate(headers):
            c = t.rows[0].cells[i].paragraphs[0].add_run(h); c.bold = True; c.font.size = Pt(9.5)
        for row in rows:
            cells = t.add_row().cells
            for i, v in enumerate(row):
                rp = cells[i].paragraphs[0].add_run(str(v)); rp.font.size = Pt(9)
        doc.add_paragraph()

    for _ in range(5):
        doc.add_paragraph()
    c = doc.add_paragraph(); c.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = c.add_run("MaestroAI"); r.bold = True; r.font.size = Pt(46); r.font.color.rgb = VIOLET
    c2 = doc.add_paragraph(); c2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = c2.add_run("Documentación técnica y funcional — Manual completo"); r.font.size = Pt(19); r.font.color.rgb = SLATE
    c3 = doc.add_paragraph(); c3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = c3.add_run("Capacidades, casos de uso, pasos y guías de configuración"); r.font.size = Pt(12); r.font.color.rgb = GREY
    c4 = doc.add_paragraph(); c4.alignment = WD_ALIGN_PARAGRAPH.CENTER
    c4.add_run("Versión 2026 · Documento confidencial").italic = True
    doc.add_page_break()

    doc.add_heading("1. Introducción y propuesta de valor", level=1)
    para("MaestroAI es una capa privada y gobernada sobre la IA para empresas de LATAM. A diferencia de un "
         "asistente genérico, cada dato que se procesa pasa por el Enrutador de Privacidad: se clasifica, se "
         "detecta PII, se minimiza y redacta, y se decide automáticamente la ruta del modelo o se bloquea — "
         "todo auditado. Sobre esa base corren capacidades verticales (RAG, casos, agentes, notebooks, "
         "automatizaciones) e integraciones con las herramientas de la empresa.")
    para("Cómo leer este documento", bold=True)
    para("La sección 4 describe cada capacidad con cuatro vistas: qué hace (Capacidad), qué se logra "
         "(Logro/valor), un ejemplo (Caso de uso) y el cómo (Pasos a seguir). La sección 6 contiene las guías "
         "técnicas de configuración paso a paso. Los apéndices listan endpoints, modelos y acciones.")

    doc.add_heading("2. Arquitectura general", level=1)
    steps([
        "Portal (Next.js 15 + Tailwind) en Vercel; instalable como PWA.",
        "API (FastAPI + SQLModel) en Render (Docker); Postgres (Supabase) en producción.",
        "RAG con almacén vectorial in-process (o Qdrant); embeddings y fragmentos cifrados.",
        "Modelos: local (Ollama), VPC (vLLM/TGI), abierto (NaN), premium (OpenAI/Claude/Gemini).",
        "Workflows con n8n gestionado; conectores y webhooks para el ecosistema.",
        "CI/CD con GitHub Actions y promoción dev → qa → main.",
    ], numbered=False)
    doc.add_heading("Flujo de una petición", level=2)
    steps([
        "El portal envía el prompt + selección de contexto a la API.",
        "El RAG recupera fragmentos por embeddings y, si está activo, los reordena (rerank).",
        "El router clasifica sensibilidad + PII, minimiza y redacta, y decide la ruta.",
        "Se genera la respuesta (con fallback seguro si el proveedor falla).",
        "Cascada opcional: refina con premium a demanda o si la respuesta es insuficiente.",
        "Se persiste el mensaje, se cita la fuente y se registra el evento de auditoría.",
    ])

    doc.add_heading("3. Modelo de privacidad y enrutamiento", level=1)
    table(["Situación del dato", "Ruta", "Notas"], [
        ["Restringido", "Local", "Nunca sale de la infraestructura."],
        ["Confidencial", "VPC (o Local)", "VPC privada con auditoría; local si no hay VPC."],
        ["PII sin alta sensibilidad", "VPC / Local", "No sale a externo por defecto."],
        ["Interno / Público", "Open (NaN)", "Empieza en NaN; premium solo a demanda/insuficiencia."],
        ["Inyección / exfiltración", "Bloqueado", "Se audita el intento."],
    ])

    doc.add_heading("4. Capacidades, logros, casos de uso y pasos", level=1)
    para("Cada capacidad se describe en cuatro vistas para separar el «qué/por qué» del «cómo».", italic=True, color=GREY)
    for cap in CAPS:
        doc.add_heading(cap["t"], level=2)
        label("Capacidad", cap["cap"])
        label("Logro / valor", cap["logro"])
        label("Caso de uso", cap["caso"])
        para("Pasos a seguir:", bold=True, color=VIOLET)
        steps(cap["pasos"])

    doc.add_heading("5. Modelos y proveedores", level=1)
    para("Todas las rutas usan endpoints compatibles con OpenAI. El proveedor abierto por defecto es NaN "
         "Builders (https://api.nan.builders/v1). Modelos publicados de NaN:")
    table(["Modelo", "Tipo", "Contexto", "Capacidades"], [list(m) for m in NAN_MODELS])
    para("Nota: la API de NaN no expone generación de imágenes (es función de su web); la sección de imágenes "
         "usa la ruta estándar OpenAI con cualquier proveedor que la exponga.", italic=True, color=GREY)

    doc.add_heading("6. Guías técnicas de configuración (paso a paso)", level=1)
    para("Cubren exactamente lo que hay que configurar para activar cada capacidad. Las variables van en "
         "Render (backend) salvo NEXT_PUBLIC_API_BASE_URL (Vercel).", italic=True, color=GREY)
    for g in GUIDES:
        doc.add_heading(g["t"], level=2)
        para(g["intro"])
        para("Pasos:", bold=True, color=VIOLET)
        steps(g["pasos"])
        if g["env"]:
            para("Variables:", bold=True, color=VIOLET)
            code(g["env"])
        if g.get("nota"):
            para("Nota: " + g["nota"], italic=True, color=GREY)

    doc.add_heading("7. Roles, permisos y multi-tenant", level=1)
    table(["Rol", "Puede"], [
        ["super_admin", "Todo, a través de todos los tenants."],
        ["admin", "Configuración del tenant: marca, facturación, usuarios, modelos, integraciones."],
        ["user", "Usar casos, agentes, documentos, automatizaciones."],
        ["security", "Auditoría, políticas y revisión de eventos."],
        ["devops", "Conectores e integraciones técnicas (fuentes de datos)."],
    ])
    para("Permisos jerárquicos por área y licencia: el rol y el área deciden qué documentos y contexto ve "
         "cada quien. General/sin área es visible para todos; Admin, Security y Super Admin ven todas las áreas.")

    doc.add_heading("8. Seguridad, privacidad y cumplimiento", level=1)
    steps([
        "Cifrado en reposo AES-256-GCM por tenant (documentos, fragmentos, credenciales).",
        "Redacción de PII antes de cualquier salida externa; minimización de contexto.",
        "Auditoría total con export a SIEM (JSONL).",
        "Endurecimiento: denylist DML/CTE en fuentes de datos, escapado OData en Graph, config solo super admin.",
        "Cumplimiento de referencia: LFPDPPP, OWASP Top 10 for LLM Apps 2025, NIST AI RMF.",
    ], numbered=False)

    doc.add_heading("9. Entornos, CI/CD y operación", level=1)
    para("Promoción dev → qa → main con PR y CI obligatorio (pytest + build). Cada push a main despliega "
         "portal (Vercel) y API (Render). Migraciones aditivas e idempotentes. Cada entorno usa sus propias "
         "variables y base de datos. Ayuda in-app en español + ayuda contextual (popup) por sección.")

    doc.add_heading("10. Roadmap y pendientes", level=1)
    steps([
        "Catálogos de configuración empresarial (en definición con el cliente).",
        "TTS (kokoro) y STT (whisper) de NaN: voz y transcripción hacia el RAG.",
        "SFTP para sistemas legados (requiere paramiko).",
        "Pendientes del cliente: proteger ramas, reconectar Microsoft con scopes, configurar premium/NaN.",
    ], numbered=False)

    doc.add_page_break()
    doc.add_heading("Apéndice A · Endpoints de la API (selección de 131)", level=1)
    for grp, eps in ENDPOINTS.items():
        doc.add_heading(grp, level=3); steps(eps, numbered=False)
    doc.add_heading("Apéndice B · Catálogo de acciones del toolkit", level=1)
    table(["Acción", "Proveedor", "Escritura", "Parámetros"], [
        ["gmail.send / outlook.send", "Google / Microsoft", "Sí", "to, subject, body"],
        ["gcal/mscal.create_event", "Google / Microsoft", "Sí", "summary, start, end, location"],
        ["gsheets.append / excel.append", "Google / Microsoft", "Sí", "id, range/table, values"],
        ["teams.post", "Microsoft", "Sí", "team_id, channel_id, message"],
        ["gsheets.read / excel.read", "Google / Microsoft", "No", "id, range/worksheet"],
        ["gcal.list / mscal.list", "Google / Microsoft", "No", "days"],
        ["onedrive.list / sharepoint.search", "Microsoft", "No", "query"],
    ])
    doc.save(path)


# ===================================================================== PPTX
def build_pptx(path: str) -> None:
    prs = Presentation(); prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    def slide(title, bullets=None, subtitle=None, cover=False):
        s = prs.slides.add_slide(blank)
        if cover:
            f = s.background.fill; f.solid(); f.fore_color.rgb = PV
        tb = s.shapes.add_textbox(PInches(0.7), PInches(2.6 if cover else 0.45), PInches(12), PInches(1.2))
        p = tb.text_frame.paragraphs[0]; p.word_wrap = True
        r = p.add_run(); r.text = title; r.font.size = PPt(36 if cover else 26); r.font.bold = True
        r.font.color.rgb = PW if cover else PV
        if subtitle:
            sb = s.shapes.add_textbox(PInches(0.7), PInches(3.9 if cover else 1.4), PInches(12), PInches(1))
            sp = sb.text_frame.paragraphs[0]; sp.word_wrap = True
            rr = sp.add_run(); rr.text = subtitle; rr.font.size = PPt(15); rr.font.color.rgb = PW if cover else PG
        if bullets:
            body = s.shapes.add_textbox(PInches(0.8), PInches(1.7), PInches(11.7), PInches(5.3))
            bf = body.text_frame; bf.word_wrap = True
            for i, b in enumerate(bullets):
                pa = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
                if isinstance(b, tuple):
                    r1 = pa.add_run(); r1.text = b[0] + ": "; r1.font.bold = True; r1.font.size = PPt(14); r1.font.color.rgb = PS
                    r2 = pa.add_run(); r2.text = b[1]; r2.font.size = PPt(12); r2.font.color.rgb = PG
                else:
                    r = pa.add_run(); r.text = "• " + b; r.font.size = PPt(14); r.font.color.rgb = PS
                pa.space_after = PPt(5)
        return s

    slide("MaestroAI", subtitle="Documentación técnica y funcional · 2026", cover=True)
    slide("Propuesta de valor", [
        ("Capa privada y gobernada", "sobre cualquier modelo."),
        ("Enrutador de Privacidad", "clasifica, redacta y decide por cada dato."),
        ("RAG por área + rerank", "respuestas con fuentes y permisos."),
        ("Actúa, no solo redacta", "toolkit con aprobación humana."),
    ])
    slide("Cómo decide la ruta", [
        ("Restringido", "local, nunca sale."), ("Confidencial/PII", "VPC o local."),
        ("Interno/Público", "NaN; premium a demanda."), ("Inyección", "bloqueado y auditado."),
    ], subtitle="El usuario nunca elige el modelo.")
    for cap in CAPS:
        slide(cap["t"], [("Logro", cap["logro"]), ("Caso de uso", cap["caso"]),
                         ("Cómo", cap["pasos"][0] + (" …" if len(cap["pasos"]) > 1 else ""))])
    slide("Configuración — lo que se conecta", [
        ("Correo", "OAuth Microsoft/Google + IMAP."),
        ("Modelos", "NaN (open), premium, Ollama (local)."),
        ("Automatización", "n8n + conectores + webhooks."),
        ("Gobierno", "rerank, eficiencia, proteger ramas."),
    ], subtitle="Guías paso a paso en la sección 6 del documento.")
    slide("Seguridad y cumplimiento", [
        "Cifrado AES-256-GCM por tenant.",
        "Redacción de PII + minimización.",
        "Auditoría total + export SIEM.",
        "LFPDPPP · OWASP LLM 2025 · NIST AI RMF.",
    ])
    slide("Gracias", subtitle="MaestroAI · Documentación 2026", cover=True)

    prs.save(path)


if __name__ == "__main__":
    os.makedirs(OUT, exist_ok=True)
    dx = os.path.join(OUT, "MaestroAI-Documentacion.docx")
    px = os.path.join(OUT, "MaestroAI-Presentacion.pptx")
    build_docx(dx); build_pptx(px)
    print("OK", dx, px)

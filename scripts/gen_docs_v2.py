"""MaestroAI v2 — genera DOS documentos (Técnico + Usuario) en Word y PPT, con
casos de uso por área. Uso: OUT_DIR=docs/generados python scripts/gen_docs_v2.py
El PDF se produce aparte convirtiendo con LibreOffice (soffice)."""
from __future__ import annotations

import os

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PColor

OUT = os.environ.get("OUT_DIR", ".")
VIOLET = RGBColor(0x6D, 0x28, 0xD9); SLATE = RGBColor(0x33, 0x41, 0x55); GREY = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x15, 0x80, 0x3D); AMBER = RGBColor(0xB4, 0x53, 0x09)
PV = PColor(0x6D, 0x28, 0xD9); PW = PColor(0xFF, 0xFF, 0xFF); PS = PColor(0x33, 0x41, 0x55); PG = PColor(0x64, 0x74, 0x8B)

# ============================================================ CASOS POR ÁREA
AREAS = [
    {"area": "Ventas", "casos": [
        {"t": "Propuesta comercial en minutos",
         "obj": "Generar una propuesta lista para enviar a partir de datos mínimos.",
         "pasos": ["Casos de uso → Propuesta comercial.", "Pon cliente, servicio, monto y notas.",
                   "Revisa el borrador (cita propuestas previas del RAG) y aprueba.", "Descarga en Word/PDF."],
         "mod": "Casos de uso · RAG (categoría propuesta_comercial) · Export", "valor": "De horas a minutos, con marca y datos citados."},
        {"t": "Resumen de correo y agenda del día",
         "obj": "Empezar el día con lo importante de la bandeja y el calendario.",
         "pasos": ["Integraciones → Conectar correo (Outlook/Gmail).", "Casos de uso → Resumen de correo y agenda.",
                   "Genera el resumen (PII redactada) y descárgalo."],
         "mod": "Toolkit lecturas (correo/calendar) · Router de privacidad", "valor": "Menos tiempo en bandeja, foco en lo que vende."},
        {"t": "Responder una licitación/RFP",
         "obj": "Construir el borrador de respuesta a unas bases.",
         "pasos": ["Notebooks → nuevo; sube las bases.", "Genera FAQ y borrador de respuesta citando las bases.",
                   "Pasa el borrador a un caso de uso para el entregable final."],
         "mod": "Notebooks · RAG", "valor": "Respuestas completas y trazables a las bases."},
        {"t": "Seguimiento: correo + alta de lead en CRM",
         "obj": "Enviar el follow-up y registrar el lead sin salir de la plataforma.",
         "pasos": ["Acciones → Enviar correo (Outlook/Gmail) → aprobar.", "Integraciones → Conector HubSpot/Salesforce → enviar el lead."],
         "mod": "Toolkit de acciones (aprobación) · Conectores", "valor": "Cero doble captura; todo auditado."},
    ]},
    {"area": "Recursos Humanos", "casos": [
        {"t": "Cartas y contratos laborales",
         "obj": "Generar oferta, carta o contrato con plantilla y datos del colaborador.",
         "pasos": ["Casos de uso → Carta/Contrato.", "Llena los datos; revisa y aprueba.", "Descarga en Word/PDF."],
         "mod": "Casos de uso · Export", "valor": "Documentos consistentes y con marca."},
        {"t": "Consultas de políticas internas",
         "obj": "Responder dudas del personal con base en el reglamento e ISO.",
         "pasos": ["Documentos → sube reglamento/políticas (área RH).", "Chat con fuentes → elegir documentos → pregunta.",
                   "La respuesta cita el documento exacto."],
         "mod": "RAG por área · Chat con fuentes", "valor": "Respuestas correctas y citadas, sin abrir manuales."},
        {"t": "Criba y resumen de CVs (con privacidad)",
         "obj": "Resumir y comparar CVs respetando datos personales.",
         "pasos": ["Documentos → sube los CVs (el router los trata como confidenciales/PII).",
                   "Chat con fuentes → pide un resumen comparativo."],
         "mod": "Router de privacidad (PII) · RAG", "valor": "Agiliza reclutamiento sin exponer PII a la nube."},
        {"t": "Onboarding automatizado",
         "obj": "Disparar el checklist y la comunicación al alta de un colaborador.",
         "pasos": ["Automatizaciones → plantilla de onboarding.", "Acción: notificar / workflow n8n / correo."],
         "mod": "Automatizaciones · n8n", "valor": "Onboarding uniforme y sin olvidos."},
    ]},
    {"area": "Finanzas", "casos": [
        {"t": "Reporte financiero / análisis",
         "obj": "Generar un reporte a partir de cifras y notas.",
         "pasos": ["Integraciones → Importar CSV/BD con las cifras (al RAG).", "Casos de uso → Reporte; aprueba.",
                   "Exporta a XLSX/PDF."],
         "mod": "Fuentes de datos (BD/CSV) · Casos de uso · Export", "valor": "Reportes con datos reales y formato listo."},
        {"t": "Consultar datos de un sistema legado",
         "obj": "Preguntar sobre datos de un ERP/BD sin API moderna.",
         "pasos": ["Integraciones → Fuentes de datos → DSN + SELECT (solo lectura) → Importar.",
                   "Chat con fuentes → pregunta sobre esos datos."],
         "mod": "Conector BD de solo lectura · RAG", "valor": "Acceso a datos legados sin riesgo de escritura."},
        {"t": "Cobranza semanal automatizada",
         "obj": "Recordatorios y reporte de cobranza cada semana.",
         "pasos": ["Automatizaciones → programada (semanal).", "Acción: workflow/conector/correo."],
         "mod": "Automatizaciones (programadas)", "valor": "Proceso recurrente sin intervención."},
        {"t": "Datos sensibles siempre protegidos",
         "obj": "Garantizar que CLABE/tarjetas no salgan a la nube.",
         "pasos": ["Trabaja normal; el router detecta PII financiera y la mantiene en ruta privada.",
                   "Verifica la ruta en el banner y en Auditoría."],
         "mod": "Router de privacidad", "valor": "Cumplimiento y tranquilidad por diseño."},
    ]},
    {"area": "Operaciones", "casos": [
        {"t": "Reporte de operación / KPIs",
         "obj": "Ver y reportar el pulso operativo.",
         "pasos": ["Tableros → describe qué medir; agrega widgets en vivo.", "Casos de uso → Reporte de operación; exporta."],
         "mod": "Tableros · Casos de uso", "valor": "Decisiones con datos en vivo."},
        {"t": "Procedimientos (SOP) con IA",
         "obj": "Consultar y resumir manuales y procedimientos.",
         "pasos": ["Notebooks → sube los SOP.", "Genera una guía o pregunta con cita."],
         "mod": "Notebooks · RAG", "valor": "Conocimiento operativo accesible al instante."},
        {"t": "Alertas por evento",
         "obj": "Reaccionar cuando ocurre algo (p. ej. documento crítico).",
         "pasos": ["Automatizaciones → disparador por evento (document_uploaded) o webhook entrante.",
                   "Acción: notificar / workflow."],
         "mod": "Automatizaciones · Webhooks", "valor": "Respuesta inmediata a eventos clave."},
        {"t": "Integrar sistemas propios",
         "obj": "Conectar herramientas a la medida sin API moderna.",
         "pasos": ["Integraciones → Conector REST / Webhook / n8n.", "Prueba y conéctalo a una automatización."],
         "mod": "Conectores · n8n", "valor": "Orquestación con el ecosistema existente."},
    ]},
    {"area": "Proyectos", "casos": [
        {"t": "Generar un SOW / enunciado de trabajo",
         "obj": "Crear un SOW metodológico a partir del input comercial.",
         "pasos": ["Agentes → Proposal/SOW.", "Da el contexto; revisa y exporta a PDF/Word."],
         "mod": "Agente Proposal/SOW · Workflow SOW", "valor": "SOWs consistentes y rápidos."},
        {"t": "Briefing y cronología del proyecto",
         "obj": "Resumir documentos del proyecto en artefactos útiles.",
         "pasos": ["Notebooks → sube los documentos del proyecto.", "Genera briefing y cronología."],
         "mod": "Notebooks", "valor": "Arranque y handoff de proyecto más claros."},
        {"t": "Diagnóstico y roadmap",
         "obj": "Evaluar (p. ej. ciberseguridad) y proponer un plan.",
         "pasos": ["Agentes → Cyber Diagnostic (u otro).", "Responde el cuestionario; obtén riesgos + roadmap."],
         "mod": "Agente Cyber Diagnostic", "valor": "Diagnóstico ejecutivo accionable."},
        {"t": "Seguimiento: eventos y tareas",
         "obj": "Agendar y dar seguimiento desde la plataforma.",
         "pasos": ["Acciones → Crear evento (Calendar) → aprobar.", "Conector para crear tareas en tu herramienta."],
         "mod": "Toolkit de acciones · Conectores", "valor": "Menos saltos entre apps."},
    ]},
    {"area": "Dirección General", "casos": [
        {"t": "Resumen ejecutivo y tableros",
         "obj": "Ver uso, costo, riesgos y documentos en un vistazo.",
         "pasos": ["Resumen → panel ejecutivo.", "Tableros → costo por ruta de modelo, casos por receta, etc."],
         "mod": "Dashboard · Tableros", "valor": "Visión 360 del negocio y de la IA."},
        {"t": "Auditoría y cumplimiento",
         "obj": "Demostrar control y trazabilidad ante auditoría.",
         "pasos": ["Auditoría → filtra por riesgo/ruta/usuario.", "Export SIEM (JSONL)."],
         "mod": "Auditoría · Export SIEM", "valor": "Evidencia para cumplimiento (LFPDPPP)."},
        {"t": "Decidir con memoria de trabajos",
         "obj": "Retomar y construir sobre trabajos previos.",
         "pasos": ["Memoria → busca por tag/tema.", "Chat → activa «Usar memoria»."],
         "mod": "Memoria + tags", "valor": "Continuidad y aprovechamiento del conocimiento."},
        {"t": "Control del costo de IA",
         "obj": "Mantener el gasto bajo control sin perder calidad.",
         "pasos": ["Admin → Eficiencia: tope de gasto, condensación, rerank.", "Cascada: premium solo a demanda."],
         "mod": "Eficiencia · Cascada", "valor": "Costo predecible; premium solo cuando aporta."},
    ]},
]

# ============================================================ SCORECARD / EXTRAS / PENDIENTES
SCORE = [
    ("Privacy Model Router (PUBLIC/INTERNAL/CONFIDENTIAL/RESTRICTED)", "Hecho"),
    ("Minimizar / anonimizar / auditar", "Hecho"),
    ("MVP: login, chat+RAG, carga segura, router, dashboard auditoría/costo", "Hecho"),
    ("Agentes: Document Intelligence, Cyber, Proposal/SOW, Executive Copilot", "Hecho"),
    ("Arquitectura por capas + multi-tenant", "Hecho"),
    ("Modelo de datos (tenants, users, agents, documents, chunks, …, audit)", "Hecho"),
    ("6 workflows operativos (Ingesta, RAG, SOW, Cyber, Centro de mando, Fine-tuning)", "Hecho (n8n)"),
    ("Cifrado AES-256 en reposo + KMS por tenant; RBAC", "Hecho"),
    ("Model adapters (OpenAI-compat, Ollama, vLLM, NaN, premium) · cost meter", "Hecho"),
    ("Export DOCX/PDF · security tests (injection, exfil, PII, cross-tenant)", "Hecho"),
    ("Ingesta: extracción de texto de PDF/DOCX", "Hecho (v2)"),
    ("Ingesta: antivirus (EICAR + ClamAV opcional) + tope de tamaño", "Hecho (v2)"),
    ("Ingesta: OCR de escaneados (tesseract en la imagen)", "Hecho (v2)"),
    ("MFA (TOTP) + códigos de respaldo, exigido en login", "Hecho (v2)"),
    ("Fine-tuning LoRA: andamiaje (datasets/anonimización/gate/jobs)", "Hecho (v2)"),
    ("Fine-tuning LoRA: entrenamiento real en GPU", "Parcial (trainer externo/lab)"),
    ("Confidential computing / TEE · mTLS / private endpoints", "Pendiente (infra)"),
]

EXTRAS = [
    "Cascada NaN → premium (premium solo a demanda) + eficiencia/condensación de tokens.",
    "Reranking del RAG (Qwen3-Reranker) para más precisión.",
    "Toolkit de acciones de ESCRITURA (Gmail/Outlook, Calendar, Sheets/Excel, Teams) con aprobación humana.",
    "Multi-cuenta de correo; contenedor de documentos por área + categorías + tratamiento.",
    "Memoria + etiquetas; Notebooks; Flujogramas navegables; Generar imágenes.",
    "Conectores legados (BD de solo lectura, CSV) → RAG; «ojito» de secretos auditado.",
    "Verificación de proveedor (Probar conexión); Ayuda in-app + ayuda contextual (popup).",
    "Pipeline dev → qa → main con CI (pytest + build) y bitácora (CHANGELOG).",
]

PENDIENTES = [
    {"t": "Fine-tuning LoRA — entrenamiento real en GPU",
     "que": "Conectar el andamiaje (ya hecho: datasets/anonimización/gate/jobs) a un trainer con GPU que entrene el adapter y lo sirva.",
     "porque": "El software ya está; falta el backend de cómputo (GPU) y servir el adapter (Ollama/vLLM).",
     "como": "Trainer (servidor GPU / App NaN / webhook n8n) que recibe el JSONL, entrena PEFT/LoRA y hace callback; servir como ruta local/VPC. Ver docs/FINETUNING-SETUP.md.",
     "dep": "GPU (local o NaN) + PEFT/transformers; Ollama/vLLM para servir."},
    {"t": "Endurecimiento enterprise — TEE / mTLS / private endpoints",
     "que": "Confidential computing (TEE) para clientes regulados, mTLS entre servicios y endpoints privados.",
     "porque": "Fase 5: clientes regulados y VPC dedicada.",
     "como": "Despliegue en VPC dedicada, mTLS entre servicios, enclaves/confidential VMs; no-retención contractual.",
     "dep": "Infraestructura (nube/VPC) y contratos con proveedores."},
]

CONFIG = [
    ("Microsoft 365 (Azure)", ["Azure → App registrations → New registration (cuentas org + personales).",
     "Redirect URI Web: https://<api>/oauth/microsoft/callback.",
     "Graph Delegated: User.Read, Mail.Read, Calendars.Read, offline_access (+ escritura: Mail.Send, Calendars.ReadWrite, ChannelMessage.Send, Files.ReadWrite.All, Sites.Read.All).",
     "Crea client secret; pon MICROSOFT_* en Render; reconecta una vez."]),
    ("Google Workspace", ["Cloud Console → habilita Gmail/Calendar/Drive/Sheets API.",
     "OAuth consent + scopes (readonly + gmail.send, spreadsheets).",
     "OAuth client (Web) con redirect https://<api>/oauth/google/callback; pon GOOGLE_* en Render."]),
    ("NaN (open) y Premium", ["Admin → Modelos externos: Base URL + modelo + API key.",
     "NaN: https://api.nan.builders/v1, modelo qwen3.6. Pulsa Probar conexión.",
     "Premium: opcional, solo se usa como escalada (máxima precisión)."]),
    ("Ollama (local privado)", ["Instala Ollama y descarga un modelo (LOCAL_MODEL).",
     "Exponlo con URL pública (cloudflared en prueba; servidor HTTPS en prod).",
     "Pon LOCAL_* en Render; ALLOW_CLOUD_FALLBACK=true mientras no haya local."]),
    ("n8n / Rerank / Ramas", ["n8n: Admin → n8n (Webhook Base URL + API key) o N8N_* en Render.",
     "Rerank: Admin → Eficiencia → activar (requiere NaN).",
     "Proteger ramas: GitHub → Settings → Branches → ruleset en main y qa (PR + checks)."]),
]


# ===================================================================== helpers
def _docx_base():
    doc = Document()
    b = doc.styles["Normal"]; b.font.name = "Calibri"; b.font.size = Pt(10.5)
    for i, sz in ((1, 18), (2, 14), (3, 12)):
        st = doc.styles[f"Heading {i}"]; st.font.color.rgb = VIOLET; st.font.size = Pt(sz)
    return doc


def _cover(doc, title, subtitle, tag):
    for _ in range(5):
        doc.add_paragraph()
    for txt, sz, col in ((title, 44, VIOLET), (subtitle, 18, SLATE), (tag, 12, GREY)):
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(txt); r.bold = sz > 20; r.font.size = Pt(sz); r.font.color.rgb = col
    v = doc.add_paragraph(); v.alignment = WD_ALIGN_PARAGRAPH.CENTER
    v.add_run("Versión 2 · 2026 · Documento confidencial").italic = True
    doc.add_page_break()


def _p(doc, text, bold=False, color=SLATE, italic=False):
    p = doc.add_paragraph(); r = p.add_run(text); r.bold = bold; r.italic = italic
    r.font.color.rgb = color; r.font.size = Pt(10.5); return p


def _lbl(doc, lbl, text, lblcolor=VIOLET):
    p = doc.add_paragraph(); r = p.add_run(lbl + ": "); r.bold = True; r.font.color.rgb = lblcolor; r.font.size = Pt(10.5)
    rr = p.add_run(text); rr.font.color.rgb = SLATE; rr.font.size = Pt(10.5)


def _steps(doc, items, numbered=True):
    for it in items:
        doc.add_paragraph(style="List Number" if numbered else "List Bullet").add_run(it)


def _table(doc, headers, rows):
    t = doc.add_table(rows=1, cols=len(headers)); t.alignment = WD_TABLE_ALIGNMENT.CENTER; t.style = "Light Grid Accent 1"
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i].paragraphs[0].add_run(h); c.bold = True; c.font.size = Pt(9.5)
    for row in rows:
        cells = t.add_row().cells
        for i, v in enumerate(row):
            run = cells[i].paragraphs[0].add_run(str(v)); run.font.size = Pt(9)
            if str(v).startswith("Pendiente"):
                run.font.color.rgb = AMBER
            elif str(v).startswith("Hecho"):
                run.font.color.rgb = GREEN
    doc.add_paragraph()


# ===================================================================== TÉCNICO
def build_tecnico(path):
    doc = _docx_base()
    _cover(doc, "MaestroAI", "Documento Técnico — v2", "Arquitectura, capacidades, casos de uso por área, pendientes y configuración")

    doc.add_heading("1. Resumen y alcance de la v2", level=1)
    _p(doc, "Esta versión 2 consolida la implementación real de MaestroAI frente al blueprint original: "
            "incorpora todo lo construido de más, deja especificado lo que falta por construir y añade casos "
            "de uso por área. La novedad de ingesta de esta versión: extracción de texto de PDF y DOCX (OCR "
            "opcional para escaneados).")

    doc.add_heading("2. Blueprint → implementación (scorecard)", level=1)
    _table(doc, ["Elemento del blueprint", "Estado"], [list(s) for s in SCORE])

    doc.add_heading("3. Construido más allá del blueprint", level=1)
    _steps(doc, EXTRAS, numbered=False)

    doc.add_heading("4. Casos de uso por área (vista técnica)", level=1)
    for a in AREAS:
        doc.add_heading(a["area"], level=2)
        for c in a["casos"]:
            doc.add_heading(c["t"], level=3)
            _lbl(doc, "Objetivo", c["obj"]); _lbl(doc, "Módulo / ruta", c["mod"]); _lbl(doc, "Valor", c["valor"])
            _p(doc, "Pasos:", bold=True, color=VIOLET); _steps(doc, c["pasos"])

    doc.add_heading("5. Pendientes por construir (especificación v2)", level=1)
    for pd in PENDIENTES:
        doc.add_heading(pd["t"], level=2)
        _lbl(doc, "Qué", pd["que"]); _lbl(doc, "Por qué", pd["porque"])
        _lbl(doc, "Cómo", pd["como"]); _lbl(doc, "Dependencias", pd["dep"], lblcolor=GREY)

    doc.add_heading("6. Guías de configuración (resumen)", level=1)
    for title, steps in CONFIG:
        doc.add_heading(title, level=3); _steps(doc, steps)
    _p(doc, "El detalle completo de variables de entorno está en el manual técnico extendido "
            "(MaestroAI-Documentacion.docx) y en docs/*.md.", italic=True, color=GREY)

    doc.add_heading("7. Seguridad y cumplimiento", level=1)
    _steps(doc, [
        "Cifrado AES-256-GCM por tenant; redacción de PII; minimización de contexto.",
        "Auditoría total con export a SIEM (JSONL); permisos por área y licencia.",
        "Aislamiento multi-tenant; endurecimiento (denylist DML/CTE, escapado OData, config solo super admin).",
        "Cumplimiento de referencia: LFPDPPP, OWASP Top 10 LLM 2025, NIST AI RMF.",
    ], numbered=False)
    doc.save(path)


# ===================================================================== USUARIO
def build_usuario(path):
    doc = _docx_base()
    _cover(doc, "MaestroAI", "Manual de Usuario — v2", "Guía práctica con casos de uso por área")

    doc.add_heading("1. ¿Qué es MaestroAI?", level=1)
    _p(doc, "Es tu plataforma de IA privada y gobernada. No eliges el modelo: el sistema clasifica cada dato, "
            "protege lo sensible y responde citando tus documentos. Tú pides el resultado; la plataforma hace "
            "el trabajo y tú apruebas.")

    doc.add_heading("2. Acceso y navegación", level=1)
    _steps(doc, [
        "Entra al portal con tu correo de empresa (instalable como app: menú → Instalar MaestroAI).",
        "Menú: Casos de uso, Chat con fuentes, Documentos, Notebooks, Generar imágenes, Memoria, Acciones, Integraciones, Tableros, Automatizaciones.",
        "Cada sección tiene un botón «? Ayuda» con la guía de esa pantalla.",
    ])

    doc.add_heading("3. Lo básico: generar un entregable", level=1)
    _steps(doc, [
        "Casos de uso → elige (propuesta, carta, reporte, licitación…).",
        "Llena los datos mínimos (objetivo, notas, formato).",
        "Revisa el borrador y pulsa «Aprobar y generar».",
        "Descarga en Word, PDF, Markdown, PPTX o XLSX.",
    ])

    doc.add_heading("4. Casos de uso por área", level=1)
    _p(doc, "Ejemplos concretos para cada área. Cada caso indica el objetivo y los pasos.", italic=True, color=GREY)
    for a in AREAS:
        doc.add_heading(a["area"], level=2)
        for c in a["casos"]:
            doc.add_heading(c["t"], level=3)
            _lbl(doc, "Objetivo", c["obj"])
            _p(doc, "Pasos:", bold=True, color=VIOLET); _steps(doc, c["pasos"])
            _lbl(doc, "Para qué sirve", c["valor"], lblcolor=GREEN)

    doc.add_heading("5. Preguntas frecuentes", level=1)
    for q, ans in [
        ("¿Por qué no elijo el modelo?", "Por seguridad: el enrutador elige la ruta más segura para cada dato."),
        ("¿Mis datos sensibles salen a la nube?", "No por defecto. Restringido → local; Confidencial/PII → privado."),
        ("¿La respuesta es confiable?", "Cuando usa tus documentos, cita la fuente y el fragmento."),
        ("Subí un PDF, ¿lo entiende?", "Sí: la plataforma extrae el texto del PDF/Word al subirlo (v2)."),
        ("Salió genérico (mock).", "Falta conectar un modelo real; avisa a tu administrador."),
    ]:
        _lbl(doc, q, ans)
    doc.save(path)


# ===================================================================== PPTX
def build_pptx(path, title, subtitle, kind):
    prs = Presentation(); prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    def slide(t, bullets=None, sub=None, cover=False):
        s = prs.slides.add_slide(blank)
        if cover:
            f = s.background.fill; f.solid(); f.fore_color.rgb = PV
        tb = s.shapes.add_textbox(PInches(0.7), PInches(2.6 if cover else 0.45), PInches(12), PInches(1.2))
        p = tb.text_frame.paragraphs[0]; p.word_wrap = True
        r = p.add_run(); r.text = t; r.font.size = PPt(34 if cover else 26); r.font.bold = True
        r.font.color.rgb = PW if cover else PV
        if sub:
            sb = s.shapes.add_textbox(PInches(0.7), PInches(3.9 if cover else 1.4), PInches(12), PInches(1))
            sp = sb.text_frame.paragraphs[0]; sp.word_wrap = True
            rr = sp.add_run(); rr.text = sub; rr.font.size = PPt(15); rr.font.color.rgb = PW if cover else PG
        if bullets:
            body = s.shapes.add_textbox(PInches(0.8), PInches(1.7), PInches(11.7), PInches(5.3))
            bf = body.text_frame; bf.word_wrap = True
            for i, b in enumerate(bullets):
                pa = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
                if isinstance(b, tuple):
                    r1 = pa.add_run(); r1.text = b[0] + ": "; r1.font.bold = True; r1.font.size = PPt(14); r1.font.color.rgb = PS
                    r2 = pa.add_run(); r2.text = b[1]; r2.font.size = PPt(12); r2.font.color.rgb = PG
                else:
                    r = pa.add_run(); r.text = "• " + b; r.font.size = PPt(14); r.font.color.rgb = PS
                pa.space_after = PPt(5)

    slide(title, sub=subtitle, cover=True)
    if kind == "tecnico":
        slide("Blueprint → implementación", [(s[0][:48], s[1]) for s in SCORE[:8]])
        slide("Construido de más", EXTRAS[:6])
        for a in AREAS:
            slide("Área: " + a["area"], [(c["t"], c["obj"]) for c in a["casos"]])
        slide("Pendientes por construir", [(p["t"], p["que"]) for p in PENDIENTES])
        slide("Configuración", [(t, "; ".join(s[:1])) for t, s in CONFIG])
    else:
        slide("Cómo generar un entregable", [
            "Casos de uso → elige.", "Llena datos mínimos.", "Revisa y aprueba.", "Descarga Word/PDF/PPTX/XLSX."])
        for a in AREAS:
            slide("Para " + a["area"], [(c["t"], c["valor"]) for c in a["casos"]])
        slide("Recuerda", [
            "No eliges el modelo: el sistema protege tus datos.",
            "Las respuestas con documentos citan la fuente.",
            "Cada sección tiene «? Ayuda».",
        ])
    slide("Gracias", sub="MaestroAI · v2 · 2026", cover=True)
    prs.save(path)


# ===================================================================== PDF (reportlab)
def build_pdf(path, kind):
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
                                    ListFlowable, ListItem, PageBreak)

    ss = getSampleStyleSheet()
    H1 = ParagraphStyle("H1", parent=ss["Heading1"], textColor=colors.HexColor("#6D28D9"), fontSize=16, spaceBefore=14)
    H2 = ParagraphStyle("H2", parent=ss["Heading2"], textColor=colors.HexColor("#6D28D9"), fontSize=13)
    H3 = ParagraphStyle("H3", parent=ss["Heading3"], textColor=colors.HexColor("#334155"), fontSize=11)
    BODY = ParagraphStyle("BODY", parent=ss["BodyText"], fontSize=10, leading=14, textColor=colors.HexColor("#334155"))
    GREYS = ParagraphStyle("GREY", parent=BODY, textColor=colors.HexColor("#64748B"), fontSize=9)
    flow = []

    def h1(t): flow.append(Paragraph(t, H1))
    def h2(t): flow.append(Paragraph(t, H2))
    def h3(t): flow.append(Paragraph(t, H3))
    def body(t, style=BODY): flow.append(Paragraph(t, style))
    def lbl(l, t): flow.append(Paragraph(f"<b><font color='#6D28D9'>{l}:</font></b> {t}", BODY))
    def bul(items):
        flow.append(ListFlowable([ListItem(Paragraph(i, BODY), leftIndent=10) for i in items], bulletType="bullet"))
    def num(items):
        flow.append(ListFlowable([ListItem(Paragraph(i, BODY), leftIndent=10) for i in items], bulletType="1"))
    def tbl(headers, rows):
        data = [[Paragraph(f"<b>{h}</b>", GREYS) for h in headers]] + \
               [[Paragraph(str(c), GREYS) for c in r] for r in rows]
        t = Table(data, colWidths=[10.5 * cm, 5.5 * cm] if len(headers) == 2 else None, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EDE9FE")),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CBD5E1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ]))
        flow.append(t); flow.append(Spacer(1, 8))

    # Cover
    flow.append(Spacer(1, 6 * cm))
    flow.append(Paragraph("<b>MaestroAI</b>", ParagraphStyle("c", parent=H1, fontSize=34, alignment=1)))
    sub = "Documento Técnico — v2" if kind == "tecnico" else "Manual de Usuario — v2"
    flow.append(Paragraph(sub, ParagraphStyle("c2", parent=H2, alignment=1, textColor=colors.HexColor("#334155"))))
    flow.append(Paragraph("Versión 2 · 2026 · Documento confidencial",
                          ParagraphStyle("c3", parent=GREYS, alignment=1)))
    flow.append(PageBreak())

    if kind == "tecnico":
        h1("1. Resumen y alcance de la v2")
        body("Consolida la implementación real frente al blueprint: incorpora lo construido de más, especifica lo "
             "pendiente y añade casos de uso por área. Novedad de ingesta v2: extracción de texto de PDF/DOCX (OCR opcional).")
        h1("2. Blueprint → implementación (scorecard)")
        tbl(["Elemento del blueprint", "Estado"], SCORE)
        h1("3. Construido más allá del blueprint"); bul(EXTRAS)
        h1("4. Casos de uso por área (vista técnica)")
        for a in AREAS:
            h2(a["area"])
            for c in a["casos"]:
                h3(c["t"]); lbl("Objetivo", c["obj"]); lbl("Módulo / ruta", c["mod"]); lbl("Valor", c["valor"])
                body("<b>Pasos:</b>"); num(c["pasos"])
        h1("5. Pendientes por construir (especificación v2)")
        for pd in PENDIENTES:
            h2(pd["t"]); lbl("Qué", pd["que"]); lbl("Por qué", pd["porque"]); lbl("Cómo", pd["como"]); lbl("Dependencias", pd["dep"])
        h1("6. Guías de configuración (resumen)")
        for title, steps in CONFIG:
            h3(title); num(steps)
        h1("7. Seguridad y cumplimiento")
        bul(["Cifrado AES-256-GCM por tenant; redacción de PII; minimización de contexto.",
             "Auditoría total + export SIEM; permisos por área y licencia; multi-tenant.",
             "LFPDPPP · OWASP Top 10 LLM 2025 · NIST AI RMF."])
    else:
        h1("1. ¿Qué es MaestroAI?")
        body("Tu plataforma de IA privada y gobernada. No eliges el modelo: el sistema clasifica cada dato, protege "
             "lo sensible y responde citando tus documentos. Tú pides el resultado y apruebas.")
        h1("2. Acceso y navegación")
        num(["Entra con tu correo de empresa (instalable como app).",
             "Menú: Casos de uso, Chat, Documentos, Notebooks, Generar imágenes, Memoria, Acciones, Integraciones, Tableros, Automatizaciones.",
             "Cada sección tiene un botón «? Ayuda»."])
        h1("3. Lo básico: generar un entregable")
        num(["Casos de uso → elige.", "Llena los datos mínimos.", "Revisa y aprueba.",
             "Descarga en Word/PDF/Markdown/PPTX/XLSX."])
        h1("4. Casos de uso por área")
        for a in AREAS:
            h2(a["area"])
            for c in a["casos"]:
                h3(c["t"]); lbl("Objetivo", c["obj"]); body("<b>Pasos:</b>"); num(c["pasos"])
                lbl("Para qué sirve", c["valor"])
        h1("5. Preguntas frecuentes")
        for q, ans in [
            ("¿Por qué no elijo el modelo?", "El enrutador elige la ruta más segura para cada dato."),
            ("¿Mis datos sensibles salen a la nube?", "No por defecto. Restringido → local; Confidencial/PII → privado."),
            ("¿La respuesta es confiable?", "Cuando usa tus documentos, cita la fuente."),
            ("Subí un PDF, ¿lo entiende?", "Sí: extrae el texto del PDF/Word al subirlo (v2)."),
            ("Salió genérico (mock).", "Falta conectar un modelo real; avisa a tu administrador."),
        ]:
            lbl(q, ans)

    SimpleDocTemplate(path, pagesize=LETTER, title="MaestroAI " + sub).build(flow)


if __name__ == "__main__":
    os.makedirs(OUT, exist_ok=True)
    build_tecnico(os.path.join(OUT, "MaestroAI-Tecnico-v2.docx"))
    build_usuario(os.path.join(OUT, "MaestroAI-Usuario-v2.docx"))
    build_pptx(os.path.join(OUT, "MaestroAI-Tecnico-v2.pptx"), "MaestroAI", "Documento Técnico — v2", "tecnico")
    build_pptx(os.path.join(OUT, "MaestroAI-Usuario-v2.pptx"), "MaestroAI", "Manual de Usuario — v2", "usuario")
    build_pdf(os.path.join(OUT, "MaestroAI-Tecnico-v2.pdf"), "tecnico")
    build_pdf(os.path.join(OUT, "MaestroAI-Usuario-v2.pdf"), "usuario")
    print("OK v2 docs en", OUT)
